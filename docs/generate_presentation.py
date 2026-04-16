#!/usr/bin/env python3
"""Generate AERO professional presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette ──
BG_DARK    = RGBColor(0x00, 0x14, 0x29)   # #001429
ACCENT     = RGBColor(0x48, 0xA2, 0xCE)   # #48A2CE accent blue
GREEN      = RGBColor(0x4C, 0x88, 0x32)   # #4C8832 accent green
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY   = RGBColor(0x99, 0x99, 0x99)
DARK_ACCENT = RGBColor(0x00, 0x2B, 0x50)  # Slightly lighter than BG for contrast
TABLE_HEADER_BG = RGBColor(0x48, 0xA2, 0xCE)
TABLE_ROW_BG1 = RGBColor(0x00, 0x1E, 0x3A)
TABLE_ROW_BG2 = RGBColor(0x00, 0x25, 0x45)

# ── Dimensions ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_NAME = "Calibri"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════
# Helper functions
# ═══════════════════════════════════════════════════════════════

def set_slide_bg(slide, color=BG_DARK):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar_top(slide, color=ACCENT, height=Inches(0.06)):
    """Add thin accent bar at the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_accent_bar_bottom(slide, color=ACCENT, height=Inches(0.04)):
    """Add thin accent bar at the bottom of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_H - height, SLIDE_W, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, total=7):
    """Add slide number at bottom right."""
    txBox = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.5), Inches(1.0), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MID_GRAY
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.RIGHT


def add_title_text(slide, text, left, top, width, height, font_size=32,
                   color=WHITE, bold=True, alignment=PP_ALIGN.LEFT):
    """Add a styled title textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_body_text(slide, text, left, top, width, height, font_size=16,
                  color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                  line_spacing=1.5):
    """Add body text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(font_size * 0.4)
    return txBox


def add_bullet_list(slide, items, left, top, width, height, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT, spacing=1.3,
                    bold_items=False):
    """Add a bulleted list with accent-colored bullet markers."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Use accent dash as bullet
        run_bullet = p.add_run()
        run_bullet.text = "\u2589  "  # Block char as bullet
        run_bullet.font.size = Pt(font_size - 2)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = FONT_NAME

        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = FONT_NAME
        run_text.font.bold = bold_items

        p.space_after = Pt(font_size * 0.6)
        p.space_before = Pt(2)

    return txBox


def add_decorative_line(slide, left, top, width, color=ACCENT, height=Inches(0.03)):
    """Add a thin decorative horizontal line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_side_accent(slide, left, top, height, color=ACCENT, width=Inches(0.05)):
    """Add a vertical accent bar (left border for sections)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=DARK_ACCENT,
                     border_color=ACCENT):
    """Add a rounded rectangle container."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape


def add_footer_text(slide, text, color=MID_GRAY, font_size=10):
    """Add footer text at the bottom center."""
    txBox = slide.shapes.add_textbox(
        Inches(0.8), SLIDE_H - Inches(0.55), Inches(8), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.font.italic = True
    p.alignment = PP_ALIGN.LEFT
    return txBox


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════

slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1)

# Large accent bar at top
add_accent_bar_top(slide1, ACCENT, Inches(0.08))

# Decorative vertical accent line on left
add_side_accent(slide1, Inches(1.0), Inches(2.0), Inches(3.5), ACCENT, Inches(0.06))

# Main title
add_title_text(
    slide1, "AERO",
    Inches(1.4), Inches(1.8), Inches(10), Inches(1.2),
    font_size=54, color=WHITE, bold=True
)

# Subtitle line 1
add_title_text(
    slide1, "System Zarzadzania Operacjami Lotniczymi",
    Inches(1.4), Inches(3.0), Inches(10), Inches(0.8),
    font_size=28, color=ACCENT, bold=False
)

# Decorative line
add_decorative_line(slide1, Inches(1.4), Inches(3.9), Inches(4.0), ACCENT)

# Subtitle line 2
add_title_text(
    slide1, "Cyfryzacja inspekcji sieci przesylowej",
    Inches(1.4), Inches(4.2), Inches(10), Inches(0.7),
    font_size=18, color=LIGHT_GRAY, bold=False
)

# Bottom accent bar
add_accent_bar_bottom(slide1, ACCENT, Inches(0.08))

# Decorative circle element (top right)
circle = slide1.shapes.add_shape(
    MSO_SHAPE.OVAL, SLIDE_W - Inches(3.5), Inches(1.0), Inches(2.5), Inches(2.5)
)
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x00, 0x1E, 0x3A)
circle.line.color.rgb = ACCENT
circle.line.width = Pt(2)

# Text inside circle
txBox = slide1.shapes.add_textbox(
    SLIDE_W - Inches(3.3), Inches(1.6), Inches(2.1), Inches(1.3)
)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "AERO"
run.font.size = Pt(20)
run.font.color.rgb = ACCENT
run.font.name = FONT_NAME
run.font.bold = True

add_slide_number(slide1, 1)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Problem i Kontekst
# ═══════════════════════════════════════════════════════════════

slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2)
add_accent_bar_top(slide2, ACCENT)

# Title
add_title_text(
    slide2, "Wyzwanie",
    Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
    font_size=36, color=WHITE
)

# Decorative line under title
add_decorative_line(slide2, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT)

# Bullet points
bullets = [
    "Inspekcje linii przesylowych 400kV wymagaja koordynacji wielu zespolow",
    "Planowanie operacji lotniczych realizowane dotad manualnie (Excel, e-mail)",
    "Brak centralnego systemu sledzenia statusow i historii zmian",
    "Potrzeba walidacji bezpieczenstwa przed kazdym lotem",
]

add_bullet_list(
    slide2, bullets,
    Inches(1.0), Inches(1.6), Inches(10.5), Inches(4.0),
    font_size=20, color=LIGHT_GRAY, bullet_color=ACCENT
)

# Footer note in accent box
box = add_rounded_rect(
    slide2, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8),
    fill_color=RGBColor(0x00, 0x1E, 0x3A), border_color=ACCENT
)

txBox = slide2.shapes.add_textbox(Inches(1.1), Inches(5.6), Inches(11.0), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT

run_icon = p.add_run()
run_icon.text = "\u25B6  "
run_icon.font.size = Pt(14)
run_icon.font.color.rgb = ACCENT
run_icon.font.name = FONT_NAME

run_text = p.add_run()
run_text.text = "Operator zarzadza >14 000 km linii najwyzszych napiec w Polsce"
run_text.font.size = Pt(16)
run_text.font.color.rgb = WHITE
run_text.font.name = FONT_NAME
run_text.font.italic = True

add_accent_bar_bottom(slide2)
add_slide_number(slide2, 2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: Architektura Systemu
# ═══════════════════════════════════════════════════════════════

slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3)
add_accent_bar_top(slide3, ACCENT)

add_title_text(
    slide3, "Architektura Technologiczna",
    Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
    font_size=36, color=WHITE
)
add_decorative_line(slide3, Inches(0.8), Inches(1.15), Inches(3.5), ACCENT)

# ── Left column: Frontend ──
col_left_x = Inches(0.8)
col_right_x = Inches(7.0)
col_top = Inches(1.6)
col_w = Inches(5.5)
col_h = Inches(5.2)

# Frontend box
add_rounded_rect(slide3, col_left_x, col_top, col_w, col_h,
                 fill_color=RGBColor(0x00, 0x1E, 0x3A), border_color=ACCENT)

# Frontend header
add_title_text(
    slide3, "FRONTEND",
    col_left_x + Inches(0.3), col_top + Inches(0.2), Inches(4.5), Inches(0.5),
    font_size=18, color=ACCENT, bold=True
)
add_decorative_line(slide3, col_left_x + Inches(0.3), col_top + Inches(0.65),
                    Inches(1.5), ACCENT, Inches(0.02))

frontend_items = [
    "React 18 + TypeScript",
    "React Router v6 (SPA)",
    "TanStack React Query",
    "Tailwind CSS + Custom Design System",
    "i18next (PL/EN)",
    "Leaflet.js -- mapy operacji",
    "Playwright -- testy E2E",
]
add_bullet_list(
    slide3, frontend_items,
    col_left_x + Inches(0.3), col_top + Inches(0.8), Inches(4.8), Inches(4.0),
    font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT
)

# ── Right column: Backend ──
add_rounded_rect(slide3, col_right_x, col_top, col_w, col_h,
                 fill_color=RGBColor(0x00, 0x1E, 0x3A), border_color=GREEN)

add_title_text(
    slide3, "BACKEND",
    col_right_x + Inches(0.3), col_top + Inches(0.2), Inches(4.5), Inches(0.5),
    font_size=18, color=GREEN, bold=True
)
add_decorative_line(slide3, col_right_x + Inches(0.3), col_top + Inches(0.65),
                    Inches(1.5), GREEN, Inches(0.02))

backend_items = [
    "Python 3.12 + FastAPI",
    "SQLAlchemy 2.0 (async)",
    "PostgreSQL 16",
    "Pydantic v2 -- walidacja",
    "JWT Authentication",
    "Docker Compose",
    "GitHub Actions CI/CD",
]
add_bullet_list(
    slide3, backend_items,
    col_right_x + Inches(0.3), col_top + Inches(0.8), Inches(4.8), Inches(4.0),
    font_size=16, color=LIGHT_GRAY, bullet_color=GREEN
)

add_accent_bar_bottom(slide3)
add_slide_number(slide3, 3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Role i Uprawnienia (RBAC)
# ═══════════════════════════════════════════════════════════════

slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4)
add_accent_bar_top(slide4, ACCENT)

add_title_text(
    slide4, "System Rol i Uprawnien",
    Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
    font_size=36, color=WHITE
)
add_decorative_line(slide4, Inches(0.8), Inches(1.15), Inches(3.0), ACCENT)

# Build a styled table
table_left = Inches(0.8)
table_top = Inches(1.6)
table_w = Inches(11.7)
rows = 5
cols = 2

table_shape = slide4.shapes.add_table(rows, cols, table_left, table_top, table_w, Inches(4.6))
table = table_shape.table

# Column widths
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(8.2)

# Table data
headers = ["Rola", "Zakres uprawnien"]
data = [
    ["Osoba planujaca\n(DE/CJI)",
     "Wprowadza operacje lotnicze, proponuje daty, monitoruje status"],
    ["Osoba nadzorujaca\n(DB)",
     "Akceptuje/odrzuca operacje, ustala planowane daty, nadzoruje caly cykl"],
    ["Pilot",
     "Tworzy zlecenia na lot, raportuje realizacje, uzupelnia czasy rzeczywiste"],
    ["Administrator",
     "Zarzadza uzytkownikami, helikopterami, zalogami, ladowiskami"],
]

def style_cell(cell, text, font_size=16, font_color=LIGHT_GRAY, bg_color=None,
               bold=False, alignment=PP_ALIGN.LEFT):
    """Style a table cell."""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

# Style header row
for j, h in enumerate(headers):
    style_cell(table.cell(0, j), h, font_size=18, font_color=WHITE,
               bg_color=TABLE_HEADER_BG, bold=True, alignment=PP_ALIGN.CENTER)

# Style data rows
for i, row_data in enumerate(data):
    bg = TABLE_ROW_BG1 if i % 2 == 0 else TABLE_ROW_BG2
    style_cell(table.cell(i + 1, 0), row_data[0], font_size=16,
               font_color=ACCENT, bg_color=bg, bold=True)
    style_cell(table.cell(i + 1, 1), row_data[1], font_size=16,
               font_color=LIGHT_GRAY, bg_color=bg)

# Remove table borders for cleaner look — set thin accent borders
from pptx.oxml.ns import qn
tbl = table._tbl
tbl_pr = tbl.tblPr if tbl.tblPr is not None else tbl._add_tblPr()

add_accent_bar_bottom(slide4)
add_slide_number(slide4, 4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Cykl Zycia Operacji
# ═══════════════════════════════════════════════════════════════

slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5)
add_accent_bar_top(slide5, ACCENT)

add_title_text(
    slide5, "Przeplyw Pracy -- Operacje i Zlecenia",
    Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
    font_size=36, color=WHITE
)
add_decorative_line(slide5, Inches(0.8), Inches(1.15), Inches(4.0), ACCENT)

# ── Flow diagram: 4 status boxes with arrows ──
statuses = [
    ("1. Wprowadzone", "Osoba\nplanujaca"),
    ("2. Potwierdzone\ndo planu", "Osoba\nnadzorujaca"),
    ("3. Zaplanowane\ndo zlecenia", "Pilot"),
    ("4. Zrealizowane", "Pilot"),
]

box_w = Inches(2.5)
box_h = Inches(1.6)
start_x = Inches(0.6)
y_top = Inches(1.7)
gap = Inches(0.55)

for i, (status, actor) in enumerate(statuses):
    x = start_x + i * (box_w + gap)

    # Status box
    box_color = ACCENT if i < 3 else GREEN
    box = add_rounded_rect(
        slide5, x, y_top, box_w, box_h,
        fill_color=RGBColor(0x00, 0x1E, 0x3A), border_color=box_color
    )

    # Status text
    txBox = slide5.shapes.add_textbox(x + Inches(0.15), y_top + Inches(0.15),
                                       box_w - Inches(0.3), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = status
    p.font.size = Pt(15)
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Actor text
    txBox2 = slide5.shapes.add_textbox(x + Inches(0.15), y_top + Inches(1.0),
                                        box_w - Inches(0.3), Inches(0.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = actor
    p2.font.size = Pt(11)
    p2.font.color.rgb = box_color
    p2.font.name = FONT_NAME
    p2.font.italic = True
    p2.alignment = PP_ALIGN.CENTER

    # Arrow between boxes (except after last)
    if i < len(statuses) - 1:
        arrow_x = x + box_w + Inches(0.05)
        arrow_y = y_top + box_h / 2 - Inches(0.12)
        txArrow = slide5.shapes.add_textbox(arrow_x, arrow_y, Inches(0.45), Inches(0.3))
        tf_a = txArrow.text_frame
        p_a = tf_a.paragraphs[0]
        p_a.text = "\u25B6"
        p_a.font.size = Pt(22)
        p_a.font.color.rgb = ACCENT
        p_a.alignment = PP_ALIGN.CENTER

# ── Additional info section ──
info_top = Inches(3.8)

# Side note box
add_rounded_rect(
    slide5, Inches(0.8), info_top, Inches(5.0), Inches(0.7),
    fill_color=RGBColor(0x00, 0x1E, 0x3A), border_color=ACCENT
)
add_body_text(
    slide5, "7 statusow operacji  x  7 statusow zlecen",
    Inches(1.1), info_top + Inches(0.12), Inches(4.5), Inches(0.5),
    font_size=17, color=ACCENT, bold=True
)

# Additional bullets
extra_bullets = [
    "Pelny audit trail -- historia zmian kazdego pola",
    "Walidacja bezpieczenstwa: waga zalogi, zasieg helikoptera, waznosc licencji",
]
add_bullet_list(
    slide5, extra_bullets,
    Inches(1.0), Inches(4.8), Inches(11.0), Inches(1.8),
    font_size=17, color=LIGHT_GRAY, bullet_color=ACCENT
)

add_accent_bar_bottom(slide5)
add_slide_number(slide5, 5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Kluczowe Funkcjonalnosci
# ═══════════════════════════════════════════════════════════════

slide6 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide6)
add_accent_bar_top(slide6, ACCENT)

add_title_text(
    slide6, "Mozliwosci Systemu",
    Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
    font_size=36, color=WHITE
)
add_decorative_line(slide6, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT)

# Two columns of features with icon text
left_features = [
    ("Planowanie operacji z trasami KML", ACCENT),
    ("Wizualizacja tras na mapie Polski", ACCENT),
    ("Workflow akceptacji z potwierdzeniami", ACCENT),
    ("Dashboard z podsumowaniem operacyjnym", ACCENT),
]

right_features = [
    ("Walidacja terytorialna (bounding box PL)", GREEN),
    ("Zarzadzanie zalogami i licencjami", GREEN),
    ("Dwujezyczny interfejs (PL/EN)", GREEN),
    ("Responsywny design (mobile/tablet)", GREEN),
]

feat_top = Inches(1.7)
feat_left_x = Inches(0.8)
feat_right_x = Inches(7.0)
feat_w = Inches(5.5)
feat_item_h = Inches(1.1)

for col_x, features, border_col in [
    (feat_left_x, left_features, ACCENT),
    (feat_right_x, right_features, GREEN)
]:
    for i, (text, color) in enumerate(features):
        y = feat_top + i * feat_item_h

        # Feature card
        card = add_rounded_rect(
            slide6, col_x, y, feat_w, Inches(0.9),
            fill_color=RGBColor(0x00, 0x1E, 0x3A), border_color=border_col
        )

        # Side accent on card
        add_side_accent(slide6, col_x, y + Inches(0.1), Inches(0.7),
                       color=border_col, width=Inches(0.05))

        # Feature text
        txBox = slide6.shapes.add_textbox(
            col_x + Inches(0.3), y + Inches(0.15), feat_w - Inches(0.5), Inches(0.6)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = FONT_NAME
        p.font.bold = False

add_accent_bar_bottom(slide6)
add_slide_number(slide6, 6)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Podsumowanie
# ═══════════════════════════════════════════════════════════════

slide7 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide7)
add_accent_bar_top(slide7, ACCENT)

add_title_text(
    slide7, "Podsumowanie i Nastepne Kroki",
    Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
    font_size=36, color=WHITE
)
add_decorative_line(slide7, Inches(0.8), Inches(1.15), Inches(3.5), ACCENT)

# Status items with checkmark-style presentation
summary_items = [
    "System gotowy do wdrozenia pilotazowego",
    "40 testow E2E pokrywajacych kluczowe scenariusze",
    "CI/CD pipeline z automatycznym deploymentem",
]

roadmap_text = "Roadmap: integracja z systemami zewnetrznymi, rozszerzona analityka, modul raportowania"

# Checkmark items (green bullets for completed items)
add_bullet_list(
    slide7, summary_items,
    Inches(1.0), Inches(1.6), Inches(10.5), Inches(2.8),
    font_size=20, color=LIGHT_GRAY, bullet_color=GREEN
)

# Roadmap in accent box
add_rounded_rect(
    slide7, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.9),
    fill_color=RGBColor(0x00, 0x1E, 0x3A), border_color=ACCENT
)

txBox = slide7.shapes.add_textbox(Inches(1.1), Inches(4.15), Inches(11.0), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]

run_arrow = p.add_run()
run_arrow.text = "\u25B6  "
run_arrow.font.size = Pt(16)
run_arrow.font.color.rgb = ACCENT
run_arrow.font.name = FONT_NAME

run_road = p.add_run()
run_road.text = roadmap_text
run_road.font.size = Pt(17)
run_road.font.color.rgb = WHITE
run_road.font.name = FONT_NAME

# Footer
add_decorative_line(slide7, Inches(0.8), Inches(5.8), Inches(11.5), ACCENT, Inches(0.02))

add_title_text(
    slide7, "AERO 2026",
    Inches(0.8), Inches(5.95), Inches(11.5), Inches(0.5),
    font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER
)

add_accent_bar_bottom(slide7, ACCENT, Inches(0.08))
add_slide_number(slide7, 7)


# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════

output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AERO_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"File size: {os.path.getsize(output_path) / 1024:.1f} KB")
print(f"Slides: {len(prs.slides)}")
